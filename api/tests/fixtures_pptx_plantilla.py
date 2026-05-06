"""Plantilla PPTX mínima con shapes AEDMI_* (SPEC §12) para tests."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from core.presentacion_plantilla import (
    NAME_ANALISIS,
    NAME_FUENTE,
    NAME_IMAGEN,
    NAME_PORTADA_TITULO,
    NAME_SUBTITULO,
    NAME_TITULO,
)

# PNG 1×1 transparente
_MINI_PNG = bytes(
    [
        0x89,
        0x50,
        0x4E,
        0x47,
        0x0D,
        0x0A,
        0x1A,
        0x0A,
        0,
        0,
        0,
        0x0D,
        0x49,
        0x48,
        0x44,
        0x52,
        0,
        0,
        0,
        1,
        0,
        0,
        0,
        1,
        8,
        6,
        0,
        0,
        0,
        0x1F,
        0x15,
        0xC4,
        0x89,
        0,
        0,
        0,
        0x0A,
        0x49,
        0x44,
        0x41,
        0x54,
        0x78,
        0x9C,
        0x63,
        0,
        1,
        0,
        0,
        5,
        0,
        1,
        0x0D,
        0x0A,
        0x2D,
        0xB4,
        0,
        0,
        0,
        0,
        0x49,
        0x45,
        0x4E,
        0x44,
        0xAE,
        0x42,
        0x60,
        0x82,
    ]
)


def mini_png() -> bytes:
    return _MINI_PNG


def escribir_plantilla_corporativa_minima(destino: Path) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[6]
    s0 = prs.slides.add_slide(layout)
    for sh in list(s0.shapes):
        sh.element.getparent().remove(sh.element)
    b0 = s0.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    b0.name = NAME_PORTADA_TITULO
    b0.text_frame.paragraphs[0].text = "PORTADA"

    s1 = prs.slides.add_slide(layout)
    for sh in list(s1.shapes):
        sh.element.getparent().remove(sh.element)
    t = s1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    t.name = NAME_TITULO
    t.text_frame.paragraphs[0].text = "TIT"
    st = s1.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.35))
    st.name = NAME_SUBTITULO
    st.text_frame.paragraphs[0].text = "SUB"
    pic = s1.shapes.add_picture(
        BytesIO(_MINI_PNG), Inches(0.5), Inches(1.3), width=Inches(4)
    )
    pic.name = NAME_IMAGEN
    a = s1.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(1.3))
    a.name = NAME_ANALISIS
    a.text_frame.paragraphs[0].text = "ANA"
    f = s1.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.45))
    f.name = NAME_FUENTE
    f.text_frame.paragraphs[0].text = "FUENTE"

    prs.save(str(destino))
