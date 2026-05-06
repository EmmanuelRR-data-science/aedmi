# api/tests/test_presentacion_export.py
from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from core.presentacion_export import (
    SeccionPptx,
    construir_pptx_bytes,
    construir_pptx_lote,
    etiqueta_origen_analisis,
    resolver_texto_analisis_exportacion,
    subtitulo_contexto_para_exportacion,
)
from core.presentacion_plantilla import (
    NAME_PORTADA_TITULO,
    NAME_TITULO,
    _get_shape_by_name,
)
from tests.fixtures_pptx_plantilla import escribir_plantilla_corporativa_minima, mini_png


def test_etiqueta_origen_app_sin_texto_generico() -> None:
    assert (etiqueta_origen_analisis("app") or "").strip() == ""
    assert (etiqueta_origen_analisis("revisado") or "").strip() == "Análisis revisado"


def test_prioridad_revisado_sobre_ia() -> None:
    t, o = resolver_texto_analisis_exportacion("  Mi texto  ", "IA descartada")
    assert o == "revisado"
    assert t == "Mi texto"


def test_fallback_ia_si_revisado_vacio() -> None:
    t, o = resolver_texto_analisis_exportacion("   \n  ", "Solo IA")
    assert o == "ia"
    assert t == "Solo IA"


def test_vacio_si_ambos_ausentes() -> None:
    t, o = resolver_texto_analisis_exportacion(None, None)
    assert o == "vacio"
    assert "Sin análisis" in t


def test_subtitulo_generico_no_se_exporta() -> None:
    assert subtitulo_contexto_para_exportacion("Dashboard AEDMI") is None
    assert subtitulo_contexto_para_exportacion("  dashboard aedmi  ") is None
    assert subtitulo_contexto_para_exportacion("Dashboaard AEDMI") is None
    assert subtitulo_contexto_para_exportacion("World Bank") == "World Bank"


def test_pptx_no_vacio() -> None:
    raw = construir_pptx_bytes("T", "Cuerpo", "revisado", None)
    assert raw.startswith(b"PK")
    assert len(raw) > 2_000


def test_pptx_lote_desde_plantilla_dos_secciones(tmp_path) -> None:
    tpl = tmp_path / "corp.pptx"
    escribir_plantilla_corporativa_minima(tpl)
    png = mini_png()
    secciones = [
        SeccionPptx(
            "Indicador A",
            "Contexto A",
            "Texto análisis A",
            "ia",
            png,
            "Fuente: X",
        ),
        SeccionPptx(
            "Indicador B",
            None,
            "Solo revisado",
            "revisado",
            None,
            None,
        ),
    ]
    raw = construir_pptx_lote(
        "Presentación prueba",
        secciones,
        template_path=str(tpl),
    )
    assert raw.startswith(b"PK")
    prs = Presentation(BytesIO(raw))
    assert len(prs.slides) == 3
    port = _get_shape_by_name(prs.slides[0], NAME_PORTADA_TITULO)
    assert port is not None and "Presentación prueba" in port.text_frame.text
    s1 = _get_shape_by_name(prs.slides[1], NAME_TITULO)
    assert s1 is not None and "Indicador A" in s1.text_frame.text
    s2 = _get_shape_by_name(prs.slides[2], NAME_TITULO)
    assert s2 is not None and "Indicador B" in s2.text_frame.text


def test_pptx_ruta_plantilla_inexistente_legacy(tmp_path) -> None:
    """Sin archivo usable se mantiene generación legacy."""
    fantasma = tmp_path / "no.pptx"
    raw = construir_pptx_lote(
        "T",
        [
            SeccionPptx(
                "I",
                None,
                "C",
                "vacio",
                None,
                None,
            ),
        ],
        template_path=str(fantasma),
    )
    assert raw.startswith(b"PK")
    prs = Presentation(BytesIO(raw))
    # Legacy: portada generada + 1 contenido = 2 slides
    assert len(prs.slides) == 2
