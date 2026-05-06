# api/core/presentacion_export.py
from __future__ import annotations

import logging
import base64
import binascii
import re
import unicodedata
from io import BytesIO
from typing import Literal, NamedTuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OrigenAnalisis = Literal["revisado", "ia", "vacio", "app"]

FONT_APTOS = "Aptos"
FONT_FALLBACK_LATIN = "Calibri"
PT_DOC = 14
PT_LEYENDA = 10
PT_PORTADA = 14


class SeccionPptx(NamedTuple):
    """Un bloque: una diapositiva con gráfica + análisis + leyenda de fuente."""

    titulo: str
    subtitulo_contexto: str | None
    cuerpo_analisis: str
    origen: OrigenAnalisis
    imagen_png: bytes | None
    leyenda_fuente: str | None


# Subtítulos que no aportan contexto (marca del producto); se omiten en PPTX/Excel/Gamma.
_GENERIC_SUBTITULOS_EXPORT = frozenset({"dashboard aedmi", "dashboaard aedmi"})


def subtitulo_contexto_para_exportacion(raw: str | None) -> str | None:
    """Devuelve ``None`` si el subtítulo es genérico o vacío; si no, el texto normalizado."""
    if not raw:
        return None
    t = " ".join(str(raw).strip().split())
    if not t:
        return None
    if t.casefold() in _GENERIC_SUBTITULOS_EXPORT:
        return None
    return t


def resolver_texto_analisis_exportacion(
    analisis_revisado: str | None,
    analisis_ia: str | None,
) -> tuple[str, OrigenAnalisis]:
    """Prioridad: revisado (no vacío) → IA → mensaje por defecto."""
    r = (analisis_revisado or "").strip()
    if r:
        return r, "revisado"
    i = (analisis_ia or "").strip()
    if i:
        return i, "ia"
    return "Sin análisis disponible para este indicador.", "vacio"


def etiqueta_origen_analisis(origen: OrigenAnalisis) -> str:
    if origen == "revisado":
        return "Análisis revisado"
    if origen == "ia":
        return "Análisis asistido por IA"
    if origen == "app":
        return ""
    return "Análisis"


def decodificar_png_base64(data_url_o_base64: str | None) -> bytes | None:
    if not data_url_o_base64 or not str(data_url_o_base64).strip():
        return None
    s = str(data_url_o_base64).strip()
    if "," in s and "base64" in s.lower():
        s = s.split(",", 1)[1]
    try:
        return base64.b64decode(s, validate=True)
    except (binascii.Error, ValueError):
        return base64.b64decode(s, validate=False)


def _slug_nombre_archivo(titulo: str) -> str:
    s = re.sub(r"[^\w\d\-_.() ]+", "", titulo, flags=re.UNICODE).strip()
    s = re.sub(r"\s+", "_", s)[:80]
    return s or "aedmi_grafica"


def _pick_layout(prs: Presentation, prefer: tuple[int, ...] = (6, 5, 1, 0)) -> object:
    for idx in prefer:
        if idx < len(prs.slide_layouts):
            return prs.slide_layouts[idx]
    return prs.slide_layouts[0]


def _fondo_blanco_slide(slide) -> None:
    """Fondo blanco sin forma rectangular (evita tapar la gráfica por orden de capas)."""
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    except Exception:
        pass


def _texto_pptx_seguro(s: str | None) -> str:
    """Unicode NFC y sin caracteres de reemplazo típicos de datos mal decodificados."""
    if not s:
        return ""
    t = unicodedata.normalize("NFC", str(s))
    t = t.replace("\ufffd", "")
    return t.strip()


def _analisis_a_bullets(texto: str, *, max_items: int = 48) -> list[str]:
    t = (texto or "").strip()
    if not t:
        return ["—"]
    bloques = [b.strip() for b in re.split(r"\n\s*\n+", t) if b.strip()]
    lineas: list[str] = []
    for b in bloques:
        if "\n" in b:
            for ln in b.split("\n"):
                s = ln.strip()
                if s:
                    lineas.append(s)
        else:
            lineas.append(b)
    if len(lineas) == 1 and len(lineas[0]) > 200:
        par = lineas[0]
        partes = re.split(r"(?<=[.!?])\s+", par)
        lineas = [p.strip() for p in partes if p.strip()]
    out = lineas[:max_items]
    return out if out else [t[:500]]


def _slide_grafica_analisis(prs: Presentation, seccion: SeccionPptx) -> None:
    """Título arriba centrado → gráfica centrada ancho completo → viñetas de análisis → leyenda de fuente abajo."""
    layout = _pick_layout(prs)
    slide = prs.slides.add_slide(layout)
    _fondo_blanco_slide(slide)

    margin = Inches(0.45)
    slide_h = prs.slide_height
    slide_w = prs.slide_width
    usable_w = slide_w - 2 * margin

    y = margin
    title_h = Inches(0.42)
    tbox = slide.shapes.add_textbox(margin, y, usable_w, title_h)
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = _texto_pptx_seguro(seccion.titulo) or "—"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_APTOS
    p.font.size = Pt(PT_DOC)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)
    y += title_h + Inches(0.06)

    if seccion.subtitulo_contexto and str(seccion.subtitulo_contexto).strip():
        sub_h = Inches(0.28)
        sub = slide.shapes.add_textbox(margin, y, usable_w, sub_h)
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = _texto_pptx_seguro(seccion.subtitulo_contexto)
        sp.alignment = PP_ALIGN.CENTER
        sp.font.name = FONT_APTOS
        sp.font.size = Pt(PT_DOC - 1)
        sp.font.italic = True
        sp.font.color.rgb = RGBColor(71, 85, 105)
        y += sub_h + Inches(0.08)
    else:
        y += Inches(0.04)

    footer_h_emu = int(Inches(0.34))
    gap_chart_analisis_emu = int(Inches(0.14))
    leyenda = _texto_pptx_seguro(seccion.leyenda_fuente)
    extra_leyenda_emu = int(Inches(0.06)) if leyenda else 0
    bottom_reserve_emu = int(margin) + footer_h_emu + extra_leyenda_emu
    y_emu = int(y)
    available_emu = int(slide_h) - bottom_reserve_emu - y_emu

    min_analisis_emu = int(Inches(1.85))
    chart_zone_emu = max(int(Inches(2.2)), int(available_emu * 48 / 100))
    if chart_zone_emu + min_analisis_emu + gap_chart_analisis_emu > available_emu:
        chart_zone_emu = max(int(Inches(1.9)), available_emu - min_analisis_emu - gap_chart_analisis_emu)

    max_pic_w_emu = int(usable_w)
    max_pic_h_emu = chart_zone_emu
    pic_top_emu = y_emu
    if seccion.imagen_png:
        try:
            im = Image.open(BytesIO(seccion.imagen_png))
            w_px, h_px = im.size
            aspect = w_px / max(h_px, 1)
        except Exception:
            aspect = 4 / 3
        pic_w_emu = max_pic_w_emu
        pic_h_emu = int(pic_w_emu / aspect)
        if pic_h_emu > max_pic_h_emu:
            pic_h_emu = max_pic_h_emu
            pic_w_emu = int(pic_h_emu * aspect)
        pic_left_emu = int(margin) + (max_pic_w_emu - pic_w_emu) // 2
        pio = BytesIO(seccion.imagen_png)
        slide.shapes.add_picture(pio, pic_left_emu, pic_top_emu, width=pic_w_emu)
        y_emu = pic_top_emu + pic_h_emu + gap_chart_analisis_emu
    else:
        aviso = slide.shapes.add_textbox(margin, y, usable_w, Inches(0.45))
        aviso.text_frame.text = "(Gráfica no capturada; exporte de nuevo con la gráfica visible.)"
        aviso.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        aviso.text_frame.paragraphs[0].font.name = FONT_APTOS
        aviso.text_frame.paragraphs[0].font.size = Pt(PT_DOC - 1)
        aviso.text_frame.paragraphs[0].font.italic = True
        aviso.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)
        y_emu += int(Inches(0.5)) + gap_chart_analisis_emu

    bullets = _analisis_a_bullets(seccion.cuerpo_analisis)
    etiqueta = (etiqueta_origen_analisis(seccion.origen) or "").strip()
    analisis_top = y_emu
    analisis_h_emu = int(slide_h) - bottom_reserve_emu - analisis_top
    if analisis_h_emu < int(Inches(1.2)):
        analisis_h_emu = int(Inches(1.2))

    body = slide.shapes.add_textbox(margin, analisis_top, usable_w, analisis_h_emu)
    rtf = body.text_frame
    rtf.word_wrap = True
    rtf.clear()
    if etiqueta:
        p_hdr = rtf.paragraphs[0]
        p_hdr.text = _texto_pptx_seguro(etiqueta)
        p_hdr.alignment = PP_ALIGN.LEFT
        p_hdr.font.name = FONT_APTOS
        p_hdr.font.size = Pt(PT_DOC)
        p_hdr.font.bold = True
        p_hdr.font.color.rgb = RGBColor(15, 23, 42)
        p_hdr.space_after = Pt(6)
        for line in bullets:
            bp = rtf.add_paragraph()
            bp.text = _texto_pptx_seguro(f"• {line}")
            bp.space_after = Pt(4)
            bp.font.name = FONT_APTOS
            bp.font.size = Pt(PT_DOC)
            bp.font.color.rgb = RGBColor(15, 23, 42)
    else:
        if not bullets:
            p0 = rtf.paragraphs[0]
            p0.text = "—"
            p0.space_after = Pt(4)
            p0.font.name = FONT_APTOS
            p0.font.size = Pt(PT_DOC)
            p0.font.color.rgb = RGBColor(15, 23, 42)
        else:
            first = True
            for line in bullets:
                bp = rtf.paragraphs[0] if first else rtf.add_paragraph()
                first = False
                bp.text = _texto_pptx_seguro(f"• {line}")
                bp.space_after = Pt(4)
                bp.font.name = FONT_APTOS
                bp.font.size = Pt(PT_DOC)
                bp.font.bold = False
                bp.font.color.rgb = RGBColor(15, 23, 42)

    if leyenda:
        foot_top_emu = int(slide_h) - int(margin) - footer_h_emu
        foot = slide.shapes.add_textbox(margin, foot_top_emu, usable_w, footer_h_emu)
        ftf = foot.text_frame
        ftf.clear()
        fp = ftf.paragraphs[0]
        fp.text = leyenda
        fp.font.name = FONT_FALLBACK_LATIN
        fp.font.size = Pt(PT_LEYENDA)
        fp.font.italic = True
        fp.font.color.rgb = RGBColor(71, 85, 105)


def anadir_portada_lote(prs: Presentation, titulo: str) -> None:
    layout = _pick_layout(prs, (0, 1, 6))
    slide = prs.slides.add_slide(layout)
    _fondo_blanco_slide(slide)
    margin = Inches(0.5)
    usable_w = prs.slide_width - 2 * margin
    box = slide.shapes.add_textbox(margin, Inches(2.4), usable_w, Inches(1.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = _texto_pptx_seguro(titulo) or '—'
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_APTOS
    p.font.size = Pt(PT_PORTADA)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)


def construir_pptx_lote(
    titulo_portada: str,
    secciones: list[SeccionPptx],
    *,
    template_path: str | None = None,
) -> bytes:
    _log = logging.getLogger(__name__)
    if template_path:
        from core import presentacion_plantilla as tpl

        if tpl.pptx_template_file_usable(template_path):
            try:
                return tpl.construir_pptx_lote_desde_plantilla(
                    template_path, titulo_portada, secciones
                )
            except tpl.PlantillaPptxError:
                raise
            except Exception as e:
                _log.warning("PPTX plantilla no usable, se usa legacy: %s", e)
    prs = Presentation()
    anadir_portada_lote(prs, titulo_portada)
    for sec in secciones:
        _slide_grafica_analisis(prs, sec)
    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def construir_pptx_bytes(
    titulo_indicador: str,
    cuerpo_analisis: str,
    origen: OrigenAnalisis,
    imagen_png: bytes | None,
    leyenda_fuente: str | None = None,
    *,
    template_path: str | None = None,
) -> bytes:
    _log = logging.getLogger(__name__)
    if template_path:
        from core import presentacion_plantilla as tpl

        if tpl.pptx_template_file_usable(template_path):
            try:
                return tpl.construir_pptx_bytes_desde_plantilla(
                    template_path,
                    titulo_indicador,
                    cuerpo_analisis,
                    origen,
                    imagen_png,
                    leyenda_fuente,
                )
            except tpl.PlantillaPptxError:
                raise
            except Exception as e:
                _log.warning("PPTX plantilla no usable, se usa legacy: %s", e)
    prs = Presentation()
    sec = SeccionPptx(
        titulo_indicador,
        None,
        cuerpo_analisis,
        origen,
        imagen_png,
        leyenda_fuente,
    )
    _slide_grafica_analisis(prs, sec)
    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def nombre_archivo_pptx(titulo: str) -> str:
    return f"{_slug_nombre_archivo(titulo)}.pptx"
