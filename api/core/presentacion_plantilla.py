# api/core/presentacion_plantilla.py
"""PPTX corporativo: abre plantilla .pptx y rellena shapes nombrados (SPEC §12)."""

from __future__ import annotations

from copy import deepcopy
from io import BytesIO
from pathlib import Path
from typing import TYPE_CHECKING

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx.shapes.placeholder import PicturePlaceholder

if TYPE_CHECKING:
    from pptx.slide import Slide

    from core.presentacion_export import OrigenAnalisis

from core.presentacion_export import (
    SeccionPptx,
    _analisis_a_bullets,
    _texto_pptx_seguro,
    etiqueta_origen_analisis,
)

MAX_BULLETS_ANALISIS_PLANTILLA = 4

# Layout alineado a plantilla corporativa ajustada en taller (p. ej. proporción 13.293" × 7.5").
LAYOUT_MARGIN_X = Inches(0.45)
# Título / subtítulo: desplazados para no solapar logo del slide master (típ. esquina sup. izq.)
LAYOUT_TITLE_LEFT = Inches(1.05)
LAYOUT_TITLE_TOP = Inches(1.2)
LAYOUT_MARGIN_BOTTOM = Inches(0.36)
LAYOUT_TITLE_H = Inches(0.4)
LAYOUT_SUBTITLE_H = Inches(0.28)
LAYOUT_GAP_SM = Inches(0.06)
LAYOUT_GAP_MD = Inches(0.08)
# Caja inicial del PNG: centrada, encaje contain tras insertar (misma idea que plantilla manual § referencia usuario).
LAYOUT_CHART_BOX_W = Inches(9.6387)
LAYOUT_CHART_BOX_H = Inches(2.625)
LAYOUT_FUENTE_H = Inches(0.32)
PT_LEYENDA_FUENTE_SLIDE = 10
PT_ANALISIS_CABECERA_SLIDE = 14
PT_ANALISIS_VIÑETAS_SLIDE = 12

# Nombres exactos del contrato (Selection Pane en PowerPoint)
NAME_PORTADA_TITULO = "AEDMI_PORTADA_TITULO"
NAME_TITULO = "AEDMI_TITULO"
NAME_SUBTITULO = "AEDMI_SUBTITULO"
NAME_IMAGEN = "AEDMI_IMAGEN"
NAME_ANALISIS = "AEDMI_ANALISIS"
NAME_FUENTE = "AEDMI_FUENTE"

_NAMES_SLIDE_CONTENIDO = frozenset(
    {NAME_TITULO, NAME_SUBTITULO, NAME_IMAGEN, NAME_ANALISIS, NAME_FUENTE}
)


class PlantillaPptxError(RuntimeError):
    """Plantilla .pptx incompleta o inválida para modo corporativo."""


def pptx_template_file_usable(path: str | None) -> bool:
    """True si existe un archivo legible en ``path`` (no vacío)."""
    if not path or not str(path).strip():
        return False
    p = Path(path.strip())
    return p.is_file() and p.stat().st_size > 0


def _iter_shapes_recursive(shapes) -> list:
    """Aplana grupo(s) para buscar por nombre."""
    out: list = []
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            out.extend(_iter_shapes_recursive(sh.shapes))
        else:
            out.append(sh)
    return out


def _get_shape_by_name(slide: Slide, name: str):
    for sh in _iter_shapes_recursive(slide.shapes):
        if sh.name == name:
            return sh
    return None


def _validar_plantilla_basica(prs: Presentation) -> None:
    if len(prs.slides) < 2:
        raise PlantillaPptxError(
            "La plantilla PPTX debe tener al menos 2 diapositivas (portada + tipo)."
        )
    s0, s1 = prs.slides[0], prs.slides[1]
    if _get_shape_by_name(s0, NAME_PORTADA_TITULO) is None:
        raise PlantillaPptxError(f'Falta la forma con name="{NAME_PORTADA_TITULO}" en slide 0.')
    for nm in _NAMES_SLIDE_CONTENIDO:
        if _get_shape_by_name(s1, nm) is None:
            raise PlantillaPptxError(f'Falta la forma con name="{nm}" en slide 1 (plantilla tipo).')


def _blips_con_embed_en(subtree) -> list:
    out = []
    for el in subtree.iter(qn("a:blip")):
        embed = el.get(qn("r:embed"))
        if embed:
            out.append((el, embed))
    return out


def _duplicar_slide_con_rewire_imagenes(prs: Presentation, fuente: Slide) -> Slide:
    """Duplica ``fuente`` al final, re-enlazando partes de imagen al slide nuevo."""
    layout = fuente.slide_layout
    destino = prs.slides.add_slide(layout)
    for shp in list(destino.shapes):
        sp = shp.element
        sp.getparent().remove(sp)
    arbol_destino = destino.shapes._spTree
    for shp in fuente.shapes:
        el = deepcopy(shp.element)
        for blip, r_id in _blips_con_embed_en(el):
            try:
                parte = fuente.part.related_part(r_id)
            except KeyError:
                continue
            ct = getattr(parte, "content_type", "") or ""
            if "image" not in ct.lower():
                continue
            try:
                blob = parte.blob
            except Exception:
                continue
            _, nuevo_rid = destino.part.get_or_add_image_part(BytesIO(blob))
            blip_el = blip
            if hasattr(blip_el, "rEmbed"):
                blip_el.rEmbed = nuevo_rid
            else:
                blip_el.set(qn("r:embed"), nuevo_rid)
        arbol_destino.insert_element_before(el, "p:extLst")
    return destino


def _eliminar_slide_por_indice(prs: Presentation, index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    sld_id = sld_ids[index]
    sld_id_lst.remove(sld_id)
    prs.part.drop_rel(sld_id.rId)


def _texto_en_shape_textframe(shape, texto: str) -> None:
    """Sustituye texto conservando el primer run cuando exista (menos pérdida de fuente)."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    t = _texto_pptx_seguro(texto)
    if not tf.paragraphs:
        return
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = t
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.text = t


def _rellenar_analisis_bullets_plantilla(shape, seccion: SeccionPptx) -> bool:
    """Cabecera de origen (si aplica) + viñetas. Devuelve si hubo línea de cabecera."""
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    tf.clear()
    cab = _texto_pptx_seguro(etiqueta_origen_analisis(seccion.origen))
    items = _analisis_a_bullets(
        seccion.cuerpo_analisis,
        max_items=MAX_BULLETS_ANALISIS_PLANTILLA,
    )
    p0 = tf.paragraphs[0]
    if cab:
        p0.text = cab
        for raw in items:
            line = _texto_pptx_seguro(raw) or "—"
            p = tf.add_paragraph()
            p.text = f"\u2022 {line}"
            p.level = 0
        return True
    if not items:
        p0.text = "—"
        return False
    p0.text = f"\u2022 {_texto_pptx_seguro(items[0]) or '—'}"
    p0.level = 0
    for raw in items[1:]:
        line = _texto_pptx_seguro(raw) or "—"
        p = tf.add_paragraph()
        p.text = f"\u2022 {line}"
        p.level = 0
    return False


def _png_alta_resolucion_si_hace_falta(png: bytes, *, min_long_edge: int = 4000) -> bytes:
    """Escala el PNG si el lado largo es bajo (p. ej. captura antigua o ratio limitado)."""
    try:
        im = Image.open(BytesIO(png))
        im.load()
        w, h = im.size
    except Exception:
        return png
    long_edge = max(w, h)
    if long_edge >= min_long_edge:
        return png
    scale = min_long_edge / long_edge
    nw, nh = max(1, int(w * scale)), max(1, int(h * scale))
    try:
        resized = im.resize((nw, nh), Image.Resampling.LANCZOS)
    except Exception:
        return png
    buf = BytesIO()
    resized.save(buf, format="PNG", optimize=False, compress_level=3)
    return buf.getvalue()


def _ajustar_imagen_en_caja_proporcional(
    shape,
    caja_l: int,
    caja_t: int,
    caja_w: int,
    caja_h: int,
    img_w: int,
    img_h: int,
) -> None:
    """Encaja la imagen en la caja del marco original sin deformar (tipo *object-fit: contain*)."""
    if caja_w <= 0 or caja_h <= 0 or img_w <= 0 or img_h <= 0:
        return
    ar_img = img_w / img_h
    cand_w = caja_w
    cand_h = int(round(cand_w / ar_img))
    if cand_h > caja_h:
        cand_h = caja_h
        cand_w = int(round(cand_h * ar_img))
    nuevo_l = int(caja_l + (caja_w - cand_w) // 2)
    nuevo_t = int(caja_t + (caja_h - cand_h) // 2)
    shape.width = cand_w
    shape.height = cand_h
    shape.left = nuevo_l
    shape.top = nuevo_t


def _reemplazar_imagen_png(shape, png: bytes) -> None:
    stream_src = _png_alta_resolucion_si_hace_falta(png)
    try:
        im = Image.open(BytesIO(stream_src))
        im.load()
        iw, ih = im.size
    except Exception:
        iw, ih = 1, 1

    caja_l, caja_t = shape.left, shape.top
    caja_w, caja_h = shape.width, shape.height
    stream = BytesIO(stream_src)
    if isinstance(shape, PicturePlaceholder):
        shape.insert_picture(stream)
        _ajustar_imagen_en_caja_proporcional(shape, caja_l, caja_t, caja_w, caja_h, iw, ih)
        return
    slide_part = shape.part
    _, r_id = slide_part.get_or_add_image_part(BytesIO(stream_src))
    pic = getattr(shape, "_pic", None)
    if pic is None:
        return
    blip = pic.blipFill.blip
    if blip is not None:
        blip.rEmbed = r_id
    _ajustar_imagen_en_caja_proporcional(shape, caja_l, caja_t, caja_w, caja_h, iw, ih)


def _aplicar_tamano_fuente_analisis_plantilla(shape, *, tiene_cabecera: bool) -> None:
    """Cabecera 14 pt si existe; viñetas 12 pt."""
    if not shape.has_text_frame:
        return
    paras = shape.text_frame.paragraphs
    if not paras:
        return
    if tiene_cabecera:
        p0 = paras[0]
        p0.font.size = Pt(PT_ANALISIS_CABECERA_SLIDE)
        for r in p0.runs:
            r.font.size = Pt(PT_ANALISIS_CABECERA_SLIDE)
        body = paras[1:]
    else:
        body = paras
    for p in body:
        p.font.size = Pt(PT_ANALISIS_VIÑETAS_SLIDE)
        for r in p.runs:
            r.font.size = Pt(PT_ANALISIS_VIÑETAS_SLIDE)


def _aplicar_tamano_fuente_shape(shape, pt: int) -> None:
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        p.font.size = Pt(pt)
        for r in p.runs:
            r.font.size = Pt(pt)


def _posicion_cabeza_y_caja_grafico(
    slide: Slide,
    prs: Presentation,
    *,
    tiene_subtitulo: bool,
    tiene_imagen: bool,
) -> tuple[int, int, int]:
    """Título, subtítulo y caja horizontal para el PNG. Devuelve (mx, w_use, slide_h)."""
    sh_t = _get_shape_by_name(slide, NAME_TITULO)
    sh_s = _get_shape_by_name(slide, NAME_SUBTITULO)
    sh_img = _get_shape_by_name(slide, NAME_IMAGEN)
    if not all([sh_t, sh_s, sh_img]):
        raise PlantillaPptxError("Slide de contenido incompleto: faltan shapes AEDMI_* .")

    sw, sh = int(prs.slide_width), int(prs.slide_height)
    mx = int(LAYOUT_MARGIN_X)
    title_left = int(LAYOUT_TITLE_LEFT)
    title_top = int(LAYOUT_TITLE_TOP)
    w_use = sw - 2 * mx
    w_title = sw - title_left - mx

    y = title_top

    sh_t.left = title_left
    sh_t.top = y
    sh_t.width = w_title
    sh_t.height = int(LAYOUT_TITLE_H)
    y += sh_t.height + int(LAYOUT_GAP_SM)

    if tiene_subtitulo:
        sh_s.left = title_left
        sh_s.top = y
        sh_s.width = w_title
        sh_s.height = int(LAYOUT_SUBTITLE_H)
        y += sh_s.height + int(LAYOUT_GAP_SM)
    else:
        sh_s.left = title_left
        sh_s.top = y
        sh_s.width = w_title
        sh_s.height = int(Inches(0.02))

    chart_w = int(LAYOUT_CHART_BOX_W)
    chart_h = int(LAYOUT_CHART_BOX_H) if tiene_imagen else int(Inches(0.02))
    img_left = max(mx, (sw - chart_w) // 2)

    sh_img.left = img_left
    sh_img.top = y
    sh_img.width = chart_w
    sh_img.height = chart_h

    return mx, w_use, sh


def _posicion_fuente_y_analisis_bajo_grafico(
    slide: Slide,
    prs: Presentation,
    *,
    mx: int,
    w_use: int,
    sh_img_bottom: int,
) -> None:
    """Leyenda (fuente) y bloque de análisis bajo el borde inferior real del gráfico."""
    sh_a = _get_shape_by_name(slide, NAME_ANALISIS)
    sh_f = _get_shape_by_name(slide, NAME_FUENTE)
    if not sh_a or not sh_f:
        raise PlantillaPptxError("Slide de contenido incompleto: faltan shapes AEDMI_* .")

    sh = int(prs.slide_height)
    mb = int(LAYOUT_MARGIN_BOTTOM)
    y = sh_img_bottom + int(LAYOUT_GAP_MD)

    sh_f.left = mx
    sh_f.top = y
    sh_f.width = w_use
    sh_f.height = int(LAYOUT_FUENTE_H)
    y += sh_f.height + int(LAYOUT_GAP_MD)

    rest = sh - mb - y
    rest = max(rest, int(Inches(1.15)))
    sh_a.left = mx
    sh_a.top = y
    sh_a.width = w_use
    sh_a.height = rest


def _rellenar_slide_contenido(slide: Slide, seccion: SeccionPptx, prs: Presentation) -> None:
    sh_t = _get_shape_by_name(slide, NAME_TITULO)
    sh_s = _get_shape_by_name(slide, NAME_SUBTITULO)
    sh_img = _get_shape_by_name(slide, NAME_IMAGEN)
    sh_a = _get_shape_by_name(slide, NAME_ANALISIS)
    sh_f = _get_shape_by_name(slide, NAME_FUENTE)
    if not all([sh_t, sh_s, sh_img, sh_a, sh_f]):
        raise PlantillaPptxError("Slide de contenido incompleto: faltan shapes AEDMI_* .")
    sub = (
        _texto_pptx_seguro(seccion.subtitulo_contexto)
        if seccion.subtitulo_contexto
        else ""
    )
    tiene_subtitulo = bool(sub)
    tiene_imagen = bool(seccion.imagen_png)
    _texto_en_shape_textframe(sh_t, seccion.titulo or "—")
    _texto_en_shape_textframe(sh_s, sub)
    mx, w_use, _sh_h = _posicion_cabeza_y_caja_grafico(
        slide,
        prs,
        tiene_subtitulo=tiene_subtitulo,
        tiene_imagen=tiene_imagen,
    )
    if tiene_imagen:
        _reemplazar_imagen_png(sh_img, seccion.imagen_png)

    img_bottom = sh_img.top + sh_img.height
    _posicion_fuente_y_analisis_bajo_grafico(
        slide, prs, mx=mx, w_use=w_use, sh_img_bottom=img_bottom
    )

    ley = _texto_pptx_seguro(seccion.leyenda_fuente) if seccion.leyenda_fuente else ""
    _texto_en_shape_textframe(sh_f, ley)
    _aplicar_tamano_fuente_shape(sh_f, PT_LEYENDA_FUENTE_SLIDE)
    tiene_cab = _rellenar_analisis_bullets_plantilla(sh_a, seccion)
    _aplicar_tamano_fuente_analisis_plantilla(sh_a, tiene_cabecera=tiene_cab)


def _rellenar_portada(slide: Slide, titulo: str) -> None:
    sh = _get_shape_by_name(slide, NAME_PORTADA_TITULO)
    if sh is None:
        raise PlantillaPptxError(f'Falta shape "{NAME_PORTADA_TITULO}" en portada.')
    _texto_en_shape_textframe(sh, titulo or "—")


def construir_pptx_lote_desde_plantilla(
    ruta_plantilla: str,
    titulo_portada: str,
    secciones: list[SeccionPptx],
) -> bytes:
    prs = Presentation(ruta_plantilla)
    _validar_plantilla_basica(prs)
    plantilla_tipo = prs.slides[1]
    for _ in secciones:
        _duplicar_slide_con_rewire_imagenes(prs, plantilla_tipo)
    _eliminar_slide_por_indice(prs, 1)
    _rellenar_portada(prs.slides[0], titulo_portada)
    for i, sec in enumerate(secciones):
        _rellenar_slide_contenido(prs.slides[1 + i], sec, prs)
    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def construir_pptx_bytes_desde_plantilla(
    ruta_plantilla: str,
    titulo_indicador: str,
    cuerpo_analisis: str,
    origen: OrigenAnalisis,
    imagen_png: bytes | None,
    leyenda_fuente: str | None = None,
) -> bytes:
    sec = SeccionPptx(
        titulo_indicador,
        None,
        cuerpo_analisis,
        origen,
        imagen_png,
        leyenda_fuente,
    )
    return construir_pptx_lote_desde_plantilla(ruta_plantilla, titulo_indicador, [sec])
