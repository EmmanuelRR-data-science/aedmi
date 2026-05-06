#!/usr/bin/env python3
"""
Deriva ``assets/plantilla-aedmi-export.pptx`` lista para la API desde
``assets/plantilla-estudio-mercado-ejemplo.pptx`` (32 diapositivas → 2).

- Conserva diapositiva 0 (portada de título) y 9 (layout con imagen + textos).
- Renombra formas a AEDMI_* y añade cajas AEDMI_SUBTITULO / AEDMI_FUENTE si faltan.
- Elimina conectores LINE de **todos** los layouts de la presentación (divisores del tema sin tocar masters ni placeholders).

Uso (desde ``api/``):
  uv run python scripts/preparar_plantilla_aedmi_export.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Inches

API_ROOT = Path(__file__).resolve().parent.parent
if str(API_ROOT) not in sys.path:
    sys.path.insert(0, str(API_ROOT))

SOURCE = API_ROOT / "assets" / "plantilla-estudio-mercado-ejemplo.pptx"
DEST = API_ROOT / "assets" / "plantilla-aedmi-export.pptx"

IDX_PORTADA = 0
IDX_TIPO = 9

NAME_PORTADA = "AEDMI_PORTADA_TITULO"
NAME_TITULO = "AEDMI_TITULO"
NAME_SUB = "AEDMI_SUBTITULO"
NAME_IMG = "AEDMI_IMAGEN"
NAME_AN = "AEDMI_ANALISIS"
NAME_FU = "AEDMI_FUENTE"


def _eliminar_conectores_line_en_layouts(prs: Presentation) -> None:
    """Quita solo formas tipo línea/conector en layouts; no altera imágenes del master ni la distribución de placeholders."""
    for layout in prs.slide_layouts:
        for sh in list(layout.shapes):
            if sh.shape_type == MSO_SHAPE_TYPE.LINE:
                sp = sh.element
                sp.getparent().remove(sp)


def _eliminar_slide_por_indice(prs: Presentation, index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    sld_id = sld_ids[index]
    sld_id_lst.remove(sld_id)
    prs.part.drop_rel(sld_id.rId)


def main() -> None:
    if not SOURCE.is_file():
        raise SystemExit(f"No existe la fuente: {SOURCE}")

    prs = Presentation(SOURCE)
    n = len(prs.slides)
    keep = {IDX_PORTADA, IDX_TIPO}
    for i in range(n - 1, -1, -1):
        if i not in keep:
            _eliminar_slide_por_indice(prs, i)

    if len(prs.slides) != 2:
        raise SystemExit(f"Se esperaban 2 slides tras recorte, hay {len(prs.slides)}")

    # --- Portada: título principal del layout ---
    portada = prs.slides[0]
    titulo_ok = False
    for sh in portada.shapes:
        if sh.is_placeholder and sh.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            sh.name = NAME_PORTADA
            titulo_ok = True
            break
    if not titulo_ok:
        raise SystemExit("Portada: no hay placeholder TITLE para mapear a AEDMI_PORTADA_TITULO")

    # --- Slide tipo: imagen + cuerpos de texto ---
    slide = prs.slides[1]
    bodies: list[tuple[int, object]] = []
    picture = None
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture = sh
        elif sh.is_placeholder and sh.placeholder_format.type == PP_PLACEHOLDER.BODY:
            bodies.append((sh.placeholder_format.idx, sh))
    bodies.sort(key=lambda x: x[0])

    if picture is None:
        raise SystemExit("Slide tipo: no se encontró ninguna forma PICTURE para AEDMI_IMAGEN")
    picture.name = NAME_IMG

    if len(bodies) >= 2:
        bodies[0][1].name = NAME_TITULO
        bodies[1][1].name = NAME_AN
    elif len(bodies) == 1:
        tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.42), Inches(12.3), Inches(0.48))
        tb.name = NAME_TITULO
        tb.text_frame.paragraphs[0].text = "\u00a0"
        bodies[0][1].name = NAME_AN
    else:
        raise SystemExit("Slide tipo: se necesita al menos un BODY para análisis")

    sub = slide.shapes.add_textbox(Inches(0.45), Inches(1.0), Inches(12.3), Inches(0.36))
    sub.name = NAME_SUB
    sub.text_frame.paragraphs[0].text = "\u00a0"

    # Pie (~última franja de la slide 16:9)
    foot = slide.shapes.add_textbox(Inches(0.45), Inches(6.92), Inches(12.3), Inches(0.42))
    foot.name = NAME_FU
    foot.text_frame.paragraphs[0].text = "\u00a0"

    _eliminar_conectores_line_en_layouts(prs)

    prs.save(DEST)

    from core.presentacion_plantilla import PlantillaPptxError, _validar_plantilla_basica

    prs_check = Presentation(DEST)
    try:
        _validar_plantilla_basica(prs_check)
    except PlantillaPptxError as e:
        raise SystemExit(str(e)) from e

    print(f"OK: {DEST}")


if __name__ == "__main__":
    main()
